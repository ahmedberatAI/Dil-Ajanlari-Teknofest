#!/usr/bin/env python
"""Baslangic sunum deck'i uretir: docs/sunum.pptx (python-pptx).

Navy/teal guvenlik temasi, 14 slayt. Juri basliklari geldiginde uyarlanir.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x02, 0x80, 0x90)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x28, 0x33)
MUTED = RGBColor(0x5A, 0x63, 0x72)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _accent(slide, top, color=TEAL, left=Inches(0.6), size=Inches(0.22)):
    sh = slide.shapes.add_shape(1, left, top, size, size)  # 1 = rectangle
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def _text(slide, left, top, width, height, lines, size=16, color=DARK,
          bold=False, bullet=False, align=PP_ALIGN.LEFT, space=10, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ("•  " + line) if bullet else line
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = font
        p.alignment = align; p.space_after = Pt(space)
    return tb


def title_slide(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    _accent(s, Inches(2.55), TEAL, left=Inches(1.0), size=Inches(0.3))
    _text(s, Inches(1.0), Inches(2.9), Inches(11), Inches(1.6), [title],
          size=46, color=WHITE, bold=True, font="Georgia")
    _text(s, Inches(1.0), Inches(4.4), Inches(11), Inches(1.0), [subtitle],
          size=22, color=ICE)
    _text(s, Inches(1.0), Inches(6.4), Inches(11), Inches(0.6), [footer],
          size=14, color=ICE)


def content_slide(title, bullets, note=None):
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    _accent(s, Inches(0.62))
    _text(s, Inches(1.0), Inches(0.5), Inches(11.5), Inches(1.0), [title],
          size=34, color=NAVY, bold=True, font="Georgia")
    _text(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(4.8), bullets,
          size=18, color=DARK, bullet=True, space=14)
    if note:
        _text(s, Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.6), [note],
              size=13, color=MUTED)


def stats_slide(title, stats):
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    _text(s, Inches(1.0), Inches(0.6), Inches(11.5), Inches(1.0), [title],
          size=34, color=WHITE, bold=True, font="Georgia")
    n = len(stats)
    gap = Inches(0.4)
    cw = (EMU_W - Inches(2.0) - gap * (n - 1)) / n
    x = Inches(1.0)
    for value, label in stats:
        _text(s, x, Inches(2.6), cw, Inches(1.2), [value],
              size=54, color=TEAL if False else ICE, bold=True, align=PP_ALIGN.CENTER, font="Arial Black")
        _text(s, x, Inches(4.0), cw, Inches(1.4), [label],
              size=16, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw + gap


def main():
    title_slide(
        "DilAjanları",
        "Yerel Video Analiz ve Karar Destek Ajanı",
        "TEKNOFEST TYDA · 3. Senaryo   |   Takım: <takım adı>  ·  Üyeler: <…>",
    )
    content_slide("Problem", [
        "Savunma sanayi ve saha kameraları sürekli, yüksek hacimli video üretir.",
        "Manuel analiz hem maliyetli hem de insan hatasına açıktır.",
        "İhtiyaç: olayları hızlı, Türkçe, yerel ve güvenilir biçimde anlamlandıran karar destek.",
    ])
    content_slide("Çözümümüz", [
        "Tamamen yerel / offline çalışan, multimodal (video + metin) yapay zekâ ajanı.",
        "Zaman damgalı olay tespiti · Türkçe özet · risk değerlendirmesi · aksiyon önerileri.",
        "Operasyonel fonksiyonları dinamik tetikleme; çıktı yapılandırılmış JSON.",
        "Dış API / bulut / ücretli yazılım bağımlılığı YOK.",
    ])
    content_slide("Mimari", [
        "Görsel-Dil Modeli: Qwen3-VL-8B (FP8), vLLM ile yerel servis (OpenAI uyumlu API).",
        "Ajan: LangGraph — 5 düğümlü, hata toleranslı durum makinesi.",
        "Öz-doğrulama + mekânsal grounding + ayarlanabilir duyarlılık modları.",
        "Arayüz: Gradio (analiz + operatör asistanı sohbeti). Tamamen yerel.",
    ], note="Diyagram: docs/architecture.md  ·  Model seçimi A/B ile kanıtlandı (Qwen2.5-VL-7B → Qwen3-VL-8B)")
    content_slide("Ajan Akışı", [
        "ingest → kare örnekleme + zaman damgalı segmentleme.",
        "perceive → iki aşamalı algı + öz-doğrulama (deduce-then-verify) + konum (grounding), paralel.",
        "reason → Türkçe özet + risk + aksiyon önerileri.",
        "act → uygun mock operasyonel fonksiyonları dinamik seç ve çağır.",
        "finalize → yapılandırılmış AnalysisResult (JSON).",
    ])
    content_slide("Agentic Bileşenler", [
        "Tools: 6 mock operasyonel fonksiyon (sağlık/güvenlik/acil durdurma/kayıt…).",
        "Memory: segmentler arası olay birikimi + çok turlu diyalog hafızası.",
        "Prompt engineering: senaryoya özel Türkçe promptlar.",
        "Dinamik, model-tabanlı araç seçimi (statik kural değil).",
    ])
    content_slide("Yenilik 1 — İki Aşamalı Algı", [
        "Katı JSON olay-promptu, düşük çözünürlüklü gerçek CCTV'de modeli aşırı temkinli yapıyordu.",
        "Çözüm: önce serbest Türkçe tarif → sonra tariften olay çıkarımı.",
        "Gerçek UCF-Crime kliplerinde tespit oranını belirgin artırdı (yüksek recall).",
    ])
    content_slide("Yenilik 2 — Risk Kalibrasyonu", [
        "Model gerçek tehlikeyi tarif edip riski sistematik düşük puanlıyordu.",
        "Severity kalibrasyonu (tehdit kelimeleri) + risk tabanı guardrail'i.",
        "Sonuç: anomalilerde risk yükseldi, normal videolarda yanlış-pozitif %0 korundu.",
    ])
    content_slide("Yenilik 3 — Model Yükseltmesi & Öz-Doğrulama", [
        "Qwen2.5-VL-7B → Qwen3-VL-8B (A/B ile): kategori adlandırma %83→%100, UCF suçları %0→%100.",
        "Daha akıcı Türkçe (Qwen3 omurgası, 119 dil) — yarışmanın çekirdek deliverable'ı.",
        "Öz-doğrulama (deduce-then-verify): ajan yüksek-riskli tespiti yeniden sorgulayıp teyit eder.",
        "Mekânsal grounding: olayın karedeki konumu (bbox → 'merkez / üst-sol') — açıklanabilirlik.",
        "Ayarlanabilir duyarlılık modları (Yüksek-Duyarlılık / Dengeli / Kararlılık).",
    ])
    content_slide("Otonomi & Diyalog", [
        "Operatör asistanı: yalnızca analize dayanır (grounding), uydurmaz.",
        "İnisiyatif alır; 'olay nerede?' → konum yanıtı; belirsizlikte açıklayıcı soru sorar.",
        "Bağlam-değişimi ve prompt-injection denemelerine dayanıklı (göreve bağlı kalır).",
        "Diyalog robustluk testi: 5.00/5 (canlı demoda gösterilecek).",
    ])
    stats_slide("Ölçüm & KPI (veriye dayalı)", [
        ("%100", "Senaryo recall + kategori (yangın/düşme)"),
        ("%0→%100", "Suç adlandırma (Qwen3-VL)"),
        ("5/5", "Diyalog robustluğu"),
        ("3.2×", "Paralelleştirme hızlanması"),
    ])
    content_slide("Performans & Ölçeklenebilirlik", [
        "Segment analizinin paralelleştirilmesi → 3.2× hızlanma (116s → 36s).",
        "vLLM + CUDA graphs → ~37 token/s; düşük gecikme.",
        "Bozuk / olaysız videoların zarif yönetimi (hata toleransı).",
        "Tamamen yerel; ölçeklenebilir mimari.",
    ])
    content_slide("Açık Kaynak & Tekrar Üretilebilirlik", [
        "Apache License 2.0 · GitHub (BilisimVadisi2026 etiketi).",
        "requirements-lock.txt ile tam sürüm kilidi; adım adım kurulum.",
        "Herkese açık veri seti linki (UCF-Crime); dokümantasyon + mimari.",
        "Benchmark altyapısı: her iyileştirme veriyle kanıtlanabilir.",
    ])
    content_slide("Sonuç & Gelecek", [
        "Çalışan, ölçülmüş, uçtan uca yerel karar destek ajanı (Qwen3-VL-8B).",
        "Güçlü yan: SOTA model + öz-doğrulama + grounding + otonomi/diyalog (5/5) + kanıtlanmış metrikler.",
        "Veriye dayalı: 8+ model/teknik kombinasyonu ölçüldü, sadece işe yarayan tutuldu.",
        "Gelecek: ses analizi, çoklu kamera, temporal lokalizasyon (window-F1), gerçek-video düşme seti.",
    ])
    title_slide("Teşekkürler", "Sorularınızı bekliyoruz — canlı demo hazır.",
                "DilAjanları · TEKNOFEST TYDA 3. Senaryo")

    out = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "sunum.pptx")
    prs.save(out)
    print(f"Olusturuldu: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slayt)")


if __name__ == "__main__":
    main()
