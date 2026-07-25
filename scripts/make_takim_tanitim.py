#!/usr/bin/env python
"""VigilantAI — Takım Tanıtım Dosyası (ilk başvuru) PPTX üreteci.
TEKNOFEST Türkçe Yapay Zeka Dil Ajanları Yarışması · 3. Senaryo · Proje: DilAjanları."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- Palet (VigilantAI: koyu-lacivert + camgöbeği, güvenlik/AI) ---
NAVY_D = RGBColor(0x0B, 0x1E, 0x3B)
NAVY = RGBColor(0x12, 0x31, 0x5E)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
CYAN_D = RGBColor(0x0E, 0x74, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTE = RGBColor(0x64, 0x74, 0x8B)
SLATE1 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE2 = RGBColor(0xE2, 0xE8, 0xF0)
ICE = RGBColor(0xCA, 0xDC, 0xFC)

HEAD_F = "Cambria"
BODY_F = "Calibri"

prs = Presentation()
prs.slide_width = I(13.333)
prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = I(13.333), I(7.5)


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def _no_shadow(sh):
    sh.shadow.inherit = False


def box(s, l, t, w, h, runs, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
        font=BODY_F, anchor=MSO_ANCHOR.TOP, wrap=True, space_after=6, line=1.05, italic=False):
    tb = s.shapes.add_textbox(I(l), I(t), I(w), I(h)); tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    items = runs if isinstance(runs, list) else [(runs, {})]
    for i, (txt, o) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = o.get("align", align)
        p.space_after = Pt(o.get("space_after", space_after)); p.space_before = Pt(0)
        p.line_spacing = o.get("line", line)
        run = p.add_run(); run.text = txt
        f = run.font
        f.size = Pt(o.get("size", size)); f.bold = o.get("bold", bold); f.italic = o.get("italic", italic)
        f.name = o.get("font", font); f.color.rgb = o.get("color", color)
    return tb


def rrect(s, l, t, w, h, fill, line=None, radius=0.08):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(l), I(t), I(w), I(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    _no_shadow(sp)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def circle(s, l, t, d, fill, text=None, tcolor=WHITE, tsize=18, tbold=True, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, I(l), I(t), I(d), I(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    else:
        sp.line.fill.background()
    _no_shadow(sp)
    if text is not None:
        tf = sp.text_frame; tf.word_wrap = False
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.size = Pt(tsize); run.font.bold = tbold; run.font.color.rgb = tcolor; run.font.name = HEAD_F
    return sp


def header(s, num, title, color=NAVY, sub=None):
    circle(s, 0.7, 0.62, 0.62, CYAN, str(num), WHITE, 22)
    box(s, 1.5, 0.6, 10.9, 0.75, title, size=30, color=color, bold=True, font=HEAD_F,
        anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        box(s, 1.5, 1.32, 10.9, 0.4, sub, size=14, color=MUTE, font=BODY_F)


def name_role(s, l, t, w, name, role):
    """İsim + rol TEK satırda (tek paragraf, iki run) — box() liste=çok-paragraf olduğu için ayrı helper."""
    tb = s.shapes.add_textbox(I(l), I(t), I(w), I(0.4)); tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.line_spacing = 1.0
    r1 = p.add_run(); r1.text = name
    r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = INK; r1.font.name = HEAD_F
    r2 = p.add_run(); r2.text = "     —  " + role
    r2.font.size = Pt(12.5); r2.font.bold = True; r2.font.color.rgb = CYAN_D; r2.font.name = BODY_F


MEMBERS = [
    ("Ahmed Berat Özer", "Takım Kaptanı", "2. sınıf", "AÖ"),
    ("Süleyman Eren Kayacılar", "Üye", "3. sınıf", "SK"),
    ("Fatma Sena Gül", "Üye", "3. sınıf", "FG"),
    ("Fatma Ceren Kar", "Üye", "3. sınıf", "FK"),
]

# ============================ SLAYT 1 — KAPAK ============================
s = slide(NAVY_D)
circle(s, 10.7, -1.1, 3.4, NAVY, None)          # dekoratif faint halka
circle(s, 11.5, 5.0, 2.6, NAVY, None)
circle(s, 0.85, 0.8, 0.9, CYAN, "V", WHITE, 34)  # logo rozeti
box(s, 0.8, 2.15, 11.7, 1.3, "VigilantAI", size=66, color=WHITE, bold=True, font=HEAD_F)
box(s, 0.85, 3.4, 11.7, 0.5, "Takım Tanıtım Dosyası", size=24, color=CYAN, bold=True, font=BODY_F)
box(s, 0.85, 4.15, 11.7, 0.9, [
    ("TEKNOFEST — Türkçe Yapay Zeka Dil Ajanları Yarışması · 3. Senaryo", {"size": 16, "color": ICE, "space_after": 4}),
    ("Proje: DilAjanları — Yerel Video Analiz & Karar Destek Ajanı", {"size": 16, "color": ICE}),
])
box(s, 0.85, 6.35, 11.7, 0.5,
    "Ahmed Berat Özer  ·  Süleyman Eren Kayacılar  ·  Fatma Sena Gül  ·  Fatma Ceren Kar",
    size=13, color=RGBColor(0x9A, 0xB2, 0xD8), font=BODY_F)

# ============================ SLAYT 2 — (1) TAKIM ADI & TANITIM ============================
s = slide()
header(s, 1, "Takım Adı ve Kısa Tanıtım")
box(s, 0.7, 2.0, 7.0, 3.6, [
    ("VigilantAI", {"size": 22, "bold": True, "color": NAVY, "space_after": 8}),
    ("Ankara Üniversitesi Yapay Zeka ve Veri Mühendisliği bölümünden dört öğrencinin "
     "kurduğu, savunma ve güvenlik odaklı yapay zeka çözümleri geliştiren bir ekiptir.",
     {"size": 16, "color": INK, "line": 1.2, "space_after": 10}),
    ("Ortak alan geçmişimizi (multimodal yapay zeka, doğal dil işleme, veri mühendisliği) "
     "gerçek dünya problemlerine uygulamak için bir araya geldik. Adımız “vigilant” "
     "(teyakkuz/uyanık) kelimesinden geliyor: kameraları insan yerine, yorulmadan ve tutarlı "
     "biçimde izleyen bir yapay zeka.", {"size": 16, "color": INK, "line": 1.2}),
])
# Sağ: hızlı-bakış kartları
rrect(s, 8.1, 2.0, 4.5, 3.6, SLATE1)
stats = [("4", "üye"), ("1", "ortak bölüm"), ("%100", "alan uyumu")]
xx = 8.45
for val, lab in stats:
    box(s, xx, 2.35, 1.35, 0.9, val, size=30, color=CYAN_D, bold=True, font=HEAD_F, align=PP_ALIGN.CENTER)
    box(s, xx, 3.25, 1.35, 0.4, lab, size=12, color=MUTE, align=PP_ALIGN.CENTER)
    xx += 1.37
box(s, 8.45, 4.0, 3.85, 1.4, [
    ("Odak alanı", {"size": 13, "bold": True, "color": NAVY, "space_after": 4}),
    ("Yerel/offline çalışan, Türkçe konuşan, açıklanabilir video-analiz ve karar-destek ajanları.",
     {"size": 13, "color": INK, "line": 1.15}),
])

# ============================ SLAYT 3 — (2) KURULUŞ AMACI ============================
s = slide()
header(s, 2, "Takımın Kuruluş Amacı")
rows = [
    ("Neden bu yarışmaya katıldık?",
     "Bölümümüzde öğrendiğimiz multimodal yapay zeka ve Türkçe doğal dil işleme bilgisini, "
     "savunma sanayine değer katan somut ve yerli bir çözüme dönüştürmek istedik."),
    ("Bu senaryoda ne çözüyoruz?",
     "Güvenlik ve saha kameraları sürekli, yüksek hacimli video üretir; bunun manuel analizi "
     "hem pahalı hem insan hatasına açıktır. Biz bu yükü; olayları zaman damgasıyla tespit eden, "
     "Türkçe özet + risk değerlendirmesi + aksiyon önerisi üreten bir karar-destek ajanına devrediyoruz."),
    ("Neyi önemsiyoruz?",
     "Tamamen yerel/offline çalışma (veri gizliliği + bağımsızlık), açıklanabilirlik ve "
     "dürüst ölçüm — abartısız, kanıtlanabilir başarım."),
]
ty = 2.05
for i, (h, b) in enumerate(rows):
    circle(s, 0.7, ty + 0.05, 0.5, NAVY, str(i + 1), WHITE, 16)
    box(s, 1.42, ty, 11.1, 0.45, h, size=18, bold=True, color=NAVY, font=HEAD_F)
    box(s, 1.42, ty + 0.5, 11.1, 0.95, b, size=15, color=INK, line=1.18)
    ty += 1.62

# ============================ SLAYT 4 — (3) EĞİTİM BİLGİLERİ ============================
s = slide()
header(s, 3, "Üyelerin Eğitim Bilgileri")
positions = [(0.7, 2.05), (6.87, 2.05), (0.7, 4.6), (6.87, 4.6)]
for (name, role, grade, ini), (l, t) in zip(MEMBERS, positions):
    rrect(s, l, t, 5.76, 2.25, SLATE1)
    circle(s, l + 0.32, t + 0.35, 1.0, NAVY, ini, WHITE, 24)
    box(s, l + 1.55, t + 0.35, 4.0, 0.5, name, size=18, bold=True, color=INK, font=HEAD_F)
    box(s, l + 1.55, t + 0.9, 4.0, 0.4, role, size=13, bold=True, color=CYAN_D)
    box(s, l + 0.32, t + 1.5, 5.15, 0.6,
        "Ankara Üniversitesi, Yapay Zeka ve Veri Mühendisliği,\nLisans Öğrencisi (" + grade + ")",
        size=13, color=MUTE, line=1.15)

# ============================ SLAYT 5 — (4) YETKİNLİKLER ============================
s = slide()
header(s, 4, "Üyelerin Yetkinlikleri")
SKILLS = [
    ("Ahmed Berat Özer", ["LLM & NLP", "Agentic sistemler (LangGraph)", "Model servisleme (vLLM)", "Prompt mühendisliği", "Python"]),
    ("Süleyman Eren Kayacılar", ["Görüntü & video işleme", "Backend", "Performans optimizasyonu", "Model servisleme", "Python"]),
    ("Fatma Sena Gül", ["Frontend (Gradio)", "Arayüz / UX tasarımı", "Veri görselleştirme", "Karar-destek çıktıları", "Python"]),
    ("Fatma Ceren Kar", ["Veri mühendisliği", "Test & benchmark", "KPI / ölçümleme", "Teknik dokümantasyon", "Python"]),
]
positions = [(0.7, 2.0), (6.87, 2.0), (0.7, 4.62), (6.87, 4.62)]
for (name, skills), (l, t) in zip(SKILLS, positions):
    rrect(s, l, t, 5.76, 2.4, WHITE, line=SLATE2)
    box(s, l + 0.3, t + 0.22, 5.2, 0.45, name, size=16, bold=True, color=NAVY, font=HEAD_F)
    # yetkinlik çipleri (2 satır)
    cx, cy = l + 0.3, t + 0.8
    for j, sk in enumerate(skills):
        w = 0.30 + len(sk) * 0.088
        if cx + w > l + 5.5:
            cx = l + 0.3; cy += 0.52
        chip = rrect(s, cx, cy, w, 0.4, SLATE1, radius=0.5)
        box(s, cx, cy, w, 0.4, sk, size=11, color=CYAN_D, bold=True, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
        cx += w + 0.14

# ============================ SLAYT 6 — (5) PROJE FİKRİ ============================
s = slide()
header(s, 5, "Proje Fikri — DilAjanları")
box(s, 0.7, 1.95, 11.9, 1.25, [
    ("Tamamen yerel/offline çalışan, Türkçe bir video-analiz ve karar-destek yapay zeka ajanı.",
     {"size": 18, "bold": True, "color": NAVY, "space_after": 6, "line": 1.15}),
    ("Savunma/saha kameralarından gelen videoyu anlar; olayları zaman damgasıyla tespit eder, "
     "Türkçe özet + risk değerlendirmesi + operatöre uygulanabilir aksiyon önerileri üretir ve "
     "gerekli operasyonel fonksiyonları (sağlık/güvenlik/kayıt) tetikler.",
     {"size": 15, "color": INK, "line": 1.2}),
])
# mini pipeline chips
steps = ["Video girişi", "Olay algısı", "Akıl yürütme", "Aksiyon önerisi", "Operasyon çağrısı"]
cx = 0.7
for i, st in enumerate(steps):
    w = 0.42 + len(st) * 0.10
    rrect(s, cx, 3.5, w, 0.58, NAVY, radius=0.28)
    box(s, cx, 3.5, w, 0.58, st, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w
    if i < len(steps) - 1:
        box(s, cx, 3.5, 0.36, 0.58, "→", size=16, color=CYAN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += 0.36
# teknoloji + çıktı kartları
rrect(s, 0.7, 4.5, 5.85, 2.15, SLATE1)
box(s, 1.0, 4.72, 5.3, 1.8, [
    ("Teknoloji yığını (hepsi açık kaynak, yerel)", {"size": 14, "bold": True, "color": NAVY, "space_after": 6}),
    ("• Qwen3-VL-8B — multimodal görüntü+dil modeli", {"size": 13, "color": INK, "space_after": 3}),
    ("• vLLM — yüksek performanslı yerel servisleme", {"size": 13, "color": INK, "space_after": 3}),
    ("• LangGraph — çok adımlı ajan akışı & araç çağrımı", {"size": 13, "color": INK, "space_after": 3}),
    ("• Tek GPU · dış API/bulut yok · Apache-2.0", {"size": 13, "color": INK}),
])
rrect(s, 6.75, 4.5, 5.85, 2.15, SLATE1)
box(s, 7.05, 4.72, 5.3, 1.8, [
    ("Ürettiği çıktı (yapılandırılmış JSON)", {"size": 14, "bold": True, "color": NAVY, "space_after": 6}),
    ("• Zaman damgalı olay listesi", {"size": 13, "color": INK, "space_after": 3}),
    ("• Türkçe genel özet", {"size": 13, "color": INK, "space_after": 3}),
    ("• Risk değerlendirmesi + gerekçe", {"size": 13, "color": INK, "space_after": 3}),
    ("• Operatöre önceliklendirilmiş aksiyon önerileri", {"size": 13, "color": INK}),
])

# ============================ SLAYT 7 — (6) GÖREV DAĞILIMI ============================
s = slide()
header(s, 6, "Görev Dağılımı")
TASKS = [
    ("Ahmed Berat Özer", "Kaptan", "AÖ",
     "Takım liderliği & koordinasyon · Model araştırma ve seçimi · LLM/Agent entegrasyonu (LangGraph) · Prompt mühendisliği"),
    ("Süleyman Eren Kayacılar", "Üye", "SK",
     "Video işleme & kare örnekleme · vLLM model servisleme ve performans optimizasyonu · Backend geliştirme"),
    ("Fatma Sena Gül", "Üye", "FG",
     "Gradio arayüzü & operatör deneyimi (UX) · Karar-destek çıktı tasarımı · Demo hazırlığı ve sunum"),
    ("Fatma Ceren Kar", "Üye", "FK",
     "Test & benchmark (KPI ölçümleme) · Veri seti yönetimi · Teknik dokümantasyon ve raporlama"),
]
ty = 1.95
for name, role, ini, task in TASKS:
    rrect(s, 0.7, ty, 11.9, 1.18, SLATE1)
    circle(s, 0.95, ty + 0.24, 0.7, NAVY if role == "Üye" else CYAN, ini, WHITE, 18)
    name_role(s, 1.9, ty + 0.2, 10.5, name, "Takım Kaptanı" if role == "Kaptan" else "Üye")
    box(s, 1.9, ty + 0.64, 10.5, 0.5, task, size=13, color=MUTE, line=1.12)
    ty += 1.30

# ============================ SLAYT 8 — (7) NEDEN BAŞARILI ============================
s = slide(NAVY_D)
circle(s, 11.2, -1.0, 3.0, NAVY, None)
circle(s, 0.7, 0.62, 0.62, CYAN, "7", WHITE, 22)
box(s, 1.5, 0.6, 11.0, 0.75, "Neden Bu Takım Başarılı Olabilir?", size=30, color=WHITE, bold=True,
    font=HEAD_F, anchor=MSO_ANCHOR.MIDDLE)
PILLARS = [
    ("Alan uyumu", "Dört üye de Ankara Üniversitesi Yapay Zeka ve Veri Mühendisliği öğrencisi — problem tam da çalıştığımız alan."),
    ("Uçtan uca teknik beceri", "LLM/agent, multimodal görüntü işleme, vLLM servisleme, backend/frontend ve benchmark; tüm katmanlar ekip içinde kapsanıyor."),
    ("Net görev paylaşımı", "Örtüşmeyen, açık roller; NLP/LLM deneyimli bir kaptan tarafından koordine ediliyor."),
    ("Kanıtlı motivasyon", "Kaptanın BTK Datathon 2026 Türkçe-metin 2.'liği + TÜBİTAK 2209-A Türkçe NLP projesi; ve şimdiden çalışan uçtan uca bir prototip."),
]
pos = [(0.7, 1.75), (6.87, 1.75), (0.7, 4.35), (6.87, 4.35)]
for (h, b), (l, t) in zip(PILLARS, pos):
    rrect(s, l, t, 5.76, 2.3, NAVY)
    circle(s, l + 0.3, t + 0.3, 0.5, CYAN, None)
    box(s, l + 0.32, t + 0.3, 0.5, 0.5, "✓", size=18, color=NAVY_D, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, l + 1.0, t + 0.32, 4.5, 0.5, h, size=17, bold=True, color=WHITE, font=HEAD_F)
    box(s, l + 0.32, t + 1.02, 5.15, 1.1, b, size=13, color=ICE, line=1.2)

prs.save("docs/takim_tanitim.pptx")
print("kaydedildi: docs/takim_tanitim.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst), "slayt")
